from flask import Flask, render_template, request, send_file, jsonify, send_from_directory
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.shared import Pt, Inches
import io
import os
import json
import re as _re
import copy
import firebase_admin
from firebase_admin import credentials, firestore
from datetime import datetime
from pptx import Presentation as PPTXPresentation
from pptx.util import Inches as PptxInches

BASE_DIR_EARLY = os.path.dirname(os.path.abspath(__file__))

firebase_json = os.getenv("FIREBASE_CREDENTIALS")

if not firebase_json:
    raise ValueError("FIREBASE_CREDENTIALS não encontrada!")

cred_dict = json.loads(firebase_json)

_cred = credentials.Certificate(cred_dict)

firebase_admin.initialize_app(_cred)
db = firestore.client()

VERDANA = 'Verdana'
TAMANHO = Pt(12)


def _fmt(run, bold=False):
    run.font.name = VERDANA
    run.font.size = TAMANHO
    if bold:
        run.bold = True


def _set_spacing_15(para):
    pPr = para._p.get_or_add_pPr()
    spacing = pPr.find(qn('w:spacing'))
    if spacing is None:
        spacing = OxmlElement('w:spacing')
        pPr.append(spacing)
    spacing.set(qn('w:line'), '360')
    spacing.set(qn('w:lineRule'), 'auto')
    jc = pPr.find(qn('w:jc'))
    if jc is None:
        jc = OxmlElement('w:jc')
        pPr.append(jc)
    jc.set(qn('w:val'), 'both')


def _centrar_paragrafo(para):
    pPr = para._p.get_or_add_pPr()
    for jc_el in pPr.findall(qn('w:jc')):
        pPr.remove(jc_el)
    jc = OxmlElement('w:jc')
    jc.set(qn('w:val'), 'center')
    pPr.append(jc)


def _inserir_linha_vazia_antes(para):
    blank = OxmlElement('w:p')
    para._element.addprevious(blank)


def _inserir_linha_vazia_apos(para):
    blank = OxmlElement('w:p')
    para._element.addnext(blank)


def inserir_tabela_apos(doc, paragrafo_ref, rows):
    tbl = doc.add_table(rows=1 + len(rows), cols=3)
    tbl.style = 'Table Grid'

    cabecalhos = ['Etapa', 'Tempo Estimado', 'Tempo da Aula']
    for j, cab in enumerate(cabecalhos):
        cell = tbl.rows[0].cells[j]
        cell.text = cab
        run = cell.paragraphs[0].runs[0]
        _fmt(run, bold=True)

    for i, (etapa, tempo, acum) in enumerate(rows):
        linha = tbl.rows[i + 1]
        for cell, txt in zip(linha.cells, [etapa, tempo, acum]):
            cell.text = txt
            for run in cell.paragraphs[0].runs:
                _fmt(run)

    paragrafo_ref._element.addnext(tbl._tbl)
    return tbl


def inserir_linha_apos(paragrafo_ref, texto, estilo='Normal'):
    novo_xml = OxmlElement('w:p')
    paragrafo_ref._element.addnext(novo_xml)
    p = DocxParagraph(novo_xml, paragrafo_ref._parent)
    try:
        p.style = estilo
    except Exception:
        pass
    _set_spacing_15(p)
    if texto:
        _fmt(p.add_run(texto))
    return p


def _paragrafo_formatado(ultimo, paragrafo_ref):
    """Cria parágrafo após 'ultimo', com pai de 'paragrafo_ref', com 1,5 e justificado."""
    novo_xml = OxmlElement('w:p')
    ultimo._element.addnext(novo_xml)
    p = DocxParagraph(novo_xml, paragrafo_ref._parent)
    try:
        p.style = 'Normal'
    except Exception:
        pass
    _set_spacing_15(p)
    return p


def inserir_lista_apos(paragrafo_ref, itens, ret=False):
    ultimo = paragrafo_ref
    for item in itens:
        if item.strip():
            novo_xml = OxmlElement('w:p')
            ultimo._element.addnext(novo_xml)
            p = DocxParagraph(novo_xml, paragrafo_ref._parent)
            try:
                p.style = 'Normal'
            except Exception:
                pass
            _set_spacing_15(p)
            _fmt(p.add_run(item.strip()))
            ultimo = p
    if ret:
        return ultimo


def inserir_lista_formatada(paragrafo_ref, itens, ret=False):
    """Objetivos/ênfase: ● negrito + primeira palavra negrito, 1,5 e justificado."""
    ultimo = paragrafo_ref
    for item in itens:
        texto = item.strip()
        if not texto:
            continue
        p = _paragrafo_formatado(ultimo, paragrafo_ref)
        _fmt(p.add_run('● '), bold=True)
        partes = texto.split(' ', 1)
        _fmt(p.add_run(partes[0]), bold=True)
        if len(partes) > 1:
            _fmt(p.add_run(' ' + partes[1]))
        ultimo = p
    if ret:
        return ultimo


def inserir_lista_diretrizes(paragrafo_ref, itens, ret=False):
    """Diretrizes: ● negrito + texto normal, 1,5 e justificado."""
    ultimo = paragrafo_ref
    for item in itens:
        texto = item.strip()
        if not texto:
            continue
        p = _paragrafo_formatado(ultimo, paragrafo_ref)
        _fmt(p.add_run('● '), bold=True)
        _fmt(p.add_run(texto))
        ultimo = p
    if ret:
        return ultimo

app = Flask(__name__)
app.config['TEMPLATES_AUTO_RELOAD'] = True

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
TEMPLATE_PATH = os.path.join(BASE_DIR, 'ESTRUTURA DO PLANO DE AULA.docx')
IMAGENS_DIR = os.path.join(BASE_DIR, 'IMAGENS')


def _remover_tabela_descricao(doc):
    for tbl in doc.tables:
        if tbl.rows and tbl.rows[0].cells[0].text.strip().upper().startswith('DESCRI'):
            tbl._tbl.getparent().remove(tbl._tbl)
            return


def _remover_imagens_apos_material(doc):
    parags = list(doc.paragraphs)
    inicio = None
    fim = None
    for i, p in enumerate(parags):
        t = p.text.strip()
        if inicio is None and '3.2' in t:
            inicio = i
        elif inicio is not None and t.startswith('4.'):
            fim = i
            break
    if inicio is None:
        return
    fim = fim if fim is not None else len(parags)
    para_remover = []
    for p in parags[inicio + 1:fim]:
        if (p._p.find('.//' + qn('w:drawing')) is not None or
                p._p.find('.//' + qn('w:pict')) is not None):
            para_remover.append(p._p)
    for p_xml in para_remover:
        parent = p_xml.getparent()
        if parent is not None:
            parent.remove(p_xml)


def substituir_paragrafo(para, substituicoes):
    texto_original = para.text
    if not any(k in texto_original for k in substituicoes):
        return

    novo_texto = texto_original
    for antigo, novo in substituicoes.items():
        novo_texto = novo_texto.replace(antigo, novo)

    if para.runs:
        para.runs[0].text = novo_texto
        for run in para.runs[1:]:
            run.text = ''


@app.route('/imagens')
def listar_imagens():
    exts = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.webp'}
    imgs = sorted([f for f in os.listdir(IMAGENS_DIR)
                   if os.path.splitext(f.lower())[1] in exts])
    return jsonify(imgs)


@app.route('/imagens/<path:filename>')
def servir_imagem(filename):
    return send_from_directory(IMAGENS_DIR, filename)


@app.route('/')
def index():
    tmpl = app.jinja_env.get_template('index.html')
    print('TEMPLATE PATH:', tmpl.filename)
    return render_template('index.html')


@app.route('/gerar', methods=['POST'])
def gerar():
    data = request.form
    doc = Document(TEMPLATE_PATH)
    _remover_tabela_descricao(doc)
    _remover_imagens_apos_material(doc)

    substituicoes = {
        '(CARGO)':             data.get('cargo', '').upper(),
        '(MODULO)':            data.get('modulo', ''),
        '(UNIDADE)':           data.get('unidade', ''),
        '(HORAS AULA)':        data.get('horas_aula', ''),
        '(CARGA HORÁRIA)':     data.get('carga_horaria', '').replace(' min', ''),
        '(MODULO DESCRIÇÃO)':  data.get('modulo_descricao', ''),
        '(ATIVIDADE DESCRIÇÃO)': data.get('atividade_descricao', ''),
        '(DESCRIÇÃO DA AULA)':   data.get('descricao_aula', ''),
        '(LOCAL DA AULA)':       data.get('local_aula', ''),
        '(DIVISÃO DA TURMA)':    data.get('divisao_turma_outro', '').strip()
                                 if data.get('divisao_turma') == 'Outro'
                                 else data.get('divisao_turma', ''),
        '(GRUPO DA TURMA)':      data.get('grupo_turma', ''),
    }

    humano_itens = []

    # Professores APC — agrupa subtipos e conta cada um
    from collections import Counter
    tipos_apc = data.getlist('prof_apc_tipo')
    contagem_apc = Counter(tipos_apc)
    for tipo, qtd in contagem_apc.items():
        humano_itens.append(f'{qtd:02d} Professor(es) APC — {tipo}')

    # Demais professores e monitores
    campos_humano = [
        ('prof_damaz',    'Professor(es) DAMAZ'),
        ('prof_dciber',   'Professor(es) DCIBER'),
        ('prof_epol',     'Professor(es) Epol'),
        ('monitor_apc',   'Monitor(es) de APC'),
        ('monitor_sala',  'Monitor(es) Sala'),
        ('monitor_epol',  'Monitor(es) Epol'),
        ('monitor_damaz', 'Monitor(es) Damaz'),
        ('monitor_dciber','Monitor(es) DCIBER'),
    ]
    for campo, label in campos_humano:
        qtd = int(data.get(campo, 0) or 0)
        if qtd > 0:
            humano_itens.append(f'{qtd:02d} {label}')

    tipos_aux   = data.getlist('auxiliar_tipo')
    outros_aux  = data.getlist('auxiliar_outro')
    contagem_aux = Counter()
    outros_livres = []
    for i, tipo in enumerate(tipos_aux):
        if tipo == 'Outro':
            txt = outros_aux[i] if i < len(outros_aux) else ''
            if txt.strip():
                outros_livres.append(txt.strip())
        else:
            contagem_aux[tipo] += 1
    for tipo, qtd in contagem_aux.items():
        humano_itens.append(f'{qtd:02d} Auxiliar(es) — {tipo}')
    for txt in outros_livres:
        humano_itens.append(f'Auxiliar — {txt}')

    mat_qtds   = data.getlist('mat_qtd')
    mat_labels = data.getlist('mat_label')
    material_itens = []
    for qtd, label in zip(mat_qtds, mat_labels):
        n = int(qtd or 0)
        if n > 0:
            material_itens.append(f'{n:02d} {label}')

    objetivo_itens = [f.strip() for f in data.getlist('obj_verbo') if f.strip()]
    enfase_itens   = [f.strip() for f in data.getlist('enfase_noun') if f.strip()]

    etapa_nomes    = data.getlist('etapa_nome')
    etapa_outros   = data.getlist('etapa_nome_outro')
    etapa_tempos   = data.getlist('etapa_tempo')
    etapa_detalhes = data.getlist('etapa_detalhe')
    etapa_imgs_raw = data.getlist('etapa_img')
    estrutura_rows    = []
    estrutura_detalhes = []
    acumulado = 0
    for i, nome in enumerate(etapa_nomes):
        nome_final = etapa_outros[i].strip() if nome == '__outro__' and i < len(etapa_outros) else nome
        try:
            tempo = int(str(etapa_tempos[i]).split(',')[0].strip() or 0) if i < len(etapa_tempos) else 0
        except (ValueError, TypeError):
            tempo = 0
        acumulado += tempo
        detalhe = etapa_detalhes[i].strip() if i < len(etapa_detalhes) else ''
        img_filename = etapa_imgs_raw[i].strip() if i < len(etapa_imgs_raw) else ''
        if nome_final:
            estrutura_rows.append((nome_final, f'{tempo} min', f'{acumulado} min'))
            estrutura_detalhes.append((nome_final, detalhe, img_filename))

    verificacao_itens  = [v.strip() for v in data.getlist('verificacao')       if v.strip()]
    ressalva_itens     = [v.strip() for v in data.getlist('ressalva')          if v.strip()]
    estrategia_intro   = data.get('estrategia_intro', '').strip()
    estrategia_bullets = [v.strip() for v in data.getlist('estrategia_bullet') if v.strip()]
    prep_cenario_nao_ha = data.get('prep_cenario_nao_ha', '') == '1'
    prep_cenario_itens  = [v.strip() for v in data.getlist('prep_cenario_item') if v.strip()]

    for para in doc.paragraphs:
        substituir_paragrafo(para, substituicoes)

    for para in doc.paragraphs:
        texto = para.text
        if 'OBJETIVOS DE APRENDIZAGEM' in texto:
            ref = inserir_linha_apos(para, 'Ao final da aula, os alunos deverão ser capazes de:', 'Normal')
            ref = inserir_lista_formatada(ref, objetivo_itens, ret=True)
            if enfase_itens:
                ref = inserir_linha_apos(ref, '', 'Normal')
                enfase_titulo = inserir_linha_apos(ref, 'Ênfase:', 'Normal')
                inserir_lista_formatada(enfase_titulo, enfase_itens)
        elif '3.1 HUMANO' in texto:
            total_prof = (int(data.get('prof_apc_qtd', 0) or 0) +
                          int(data.get('prof_damaz', 0) or 0) +
                          int(data.get('prof_dciber', 0) or 0) +
                          int(data.get('prof_epol', 0) or 0))
            total_mon  = (int(data.get('monitor_apc', 0) or 0) +
                          int(data.get('monitor_sala', 0) or 0) +
                          int(data.get('monitor_epol', 0) or 0) +
                          int(data.get('monitor_damaz', 0) or 0) +
                          int(data.get('monitor_dciber', 0) or 0))
            total_aux  = int(data.get('aux_qtd', 0) or 0)
            partes_rh = []
            if total_prof > 0: partes_rh.append(f'{total_prof:02d} professores')
            if total_mon  > 0: partes_rh.append(f'{total_mon:02d} monitores')
            if total_aux  > 0: partes_rh.append(f'{total_aux:02d} auxiliares')
            if partes_rh:
                _fmt(para.add_run(f' ({", ".join(partes_rh)})'))
            inserir_lista_apos(para, humano_itens)
        elif '3.2 MATERIAL' in texto:
            inserir_lista_apos(para, material_itens)
        elif '3.3' in texto and 'CEN' in texto.upper():
            if prep_cenario_nao_ha:
                inserir_lista_apos(para, ['Não há.'])
            elif prep_cenario_itens:
                ref = para
                for i, item in enumerate(prep_cenario_itens):
                    ref = inserir_linha_apos(ref, f'3.3.{i + 1} {item}', 'Normal')
        elif 'ESTRATÉGIA DE ENSINO' in texto:
            ref = para
            if estrategia_intro:
                ref = inserir_linha_apos(ref, estrategia_intro, 'Normal')
            inserir_lista_diretrizes(ref, estrategia_bullets)
        elif 'ESTRUTURA GERAL DA AULA' in texto:
            if estrutura_rows:
                tbl = inserir_tabela_apos(doc, para, estrutura_rows)
                ultimo = tbl._tbl
                for idx, (nome_etapa, detalhe, img_filename) in enumerate(estrutura_detalhes):
                    p_blank = OxmlElement('w:p')
                    ultimo.addnext(p_blank)
                    p_tit_xml = OxmlElement('w:p')
                    p_blank.addnext(p_tit_xml)
                    p_tit = DocxParagraph(p_tit_xml, para._parent)
                    _fmt(p_tit.add_run(f'4.{idx + 1} {nome_etapa}'), bold=True)
                    p_blank_tit = OxmlElement('w:p')
                    p_tit_xml.addnext(p_blank_tit)
                    p_det_xml = OxmlElement('w:p')
                    p_blank_tit.addnext(p_det_xml)
                    p_det = DocxParagraph(p_det_xml, para._parent)
                    if detalhe:
                        _fmt(p_det.add_run(detalhe))
                    ultimo = p_det_xml
                    if img_filename:
                        img_fns = [f.strip() for f in img_filename.split('|') if f.strip()]
                        ref_xml = p_det_xml
                        for img_fn in img_fns:
                            img_path = os.path.join(IMAGENS_DIR, img_fn)
                            if os.path.isfile(img_path):
                                p_img_xml = OxmlElement('w:p')
                                ref_xml.addnext(p_img_xml)
                                p_img = DocxParagraph(p_img_xml, para._parent)
                                try:
                                    p_img.style = 'Normal'
                                except Exception:
                                    pass
                                p_img.add_run().add_picture(img_path, width=Inches(5.5))
                                _set_spacing_15(p_img)
                                ref_xml = p_img_xml
                        ultimo = ref_xml
        elif 'Ressalvas Did' in texto:
            inserir_lista_apos(para, ressalva_itens)
        elif 'VERIFICAÇÃO DE APRENDIZAGEM' in texto:
            ref = para
            for item in verificacao_itens:
                ref = inserir_linha_apos(ref, item, 'Normal')
        elif 'AVALIA' in texto.upper() and 'AULA' in texto.upper():
            para.paragraph_format.page_break_before = True

    # Aplica espaçamento 1,5 e justificado em todos os parágrafos do documento
    for para in doc.paragraphs:
        _set_spacing_15(para)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for para in cell.paragraphs:
                    _set_spacing_15(para)

    # Centraliza as 5 primeiras linhas da primeira página
    for para in doc.paragraphs[:5]:
        _centrar_paragrafo(para)

    # Adiciona linha vazia antes e após cada cabeçalho de seção
    SECAO_CHAVES = [
        'OBJETIVOS DE APRENDIZAGEM',
        '3.1 HUMANO',
        '3.2 MATERIAL',
        'ESTRATÉGIA DE ENSINO',
        'ESTRUTURA GERAL DA AULA',
        'VERIFICAÇÃO DE APRENDIZAGEM',
    ]

    def _e_cabecalho_secao(t):
        tu = t.upper()
        if any(k in tu for k in SECAO_CHAVES):
            return True
        if '3.3' in t and 'CEN' in tu:
            return True
        if 'Ressalvas Did' in t:
            return True
        return False

    cabecalhos = [p for p in doc.paragraphs if _e_cabecalho_secao(p.text)]
    for para in cabecalhos:
        _inserir_linha_vazia_antes(para)
        _inserir_linha_vazia_apos(para)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    mod = data.get('modulo', 'XX').zfill(2)
    uni = data.get('unidade', 'XX').zfill(2)
    filename = f"Plano_Aula_M{mod}_U{uni}.docx"

    return send_file(
        buf,
        download_name=filename,
        as_attachment=True,
        mimetype='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
    )


@app.route('/salvar', methods=['POST'])
def salvar():
    data = request.json
    nome = data.get('nome', '').strip() or 'Sem título'
    doc_ref = db.collection('projetos').document()
    doc_ref.set({
        'nome': nome,
        'usuario': data.get('usuario', ''),
        'dados': data.get('dados', {}),
        'status': data.get('status', 'rascunho'),
        'criado_em': datetime.now().isoformat(),
        'atualizado_em': datetime.now().isoformat(),
    })
    return jsonify({'id': doc_ref.id, 'nome': nome})


@app.route('/projetos', methods=['GET'])
def listar_projetos():
    docs = db.collection('projetos').order_by('atualizado_em', direction=firestore.Query.DESCENDING).stream()
    result = []
    for doc in docs:
        d = doc.to_dict()
        result.append({
            'id': doc.id,
            'nome': d.get('nome', ''),
            'usuario': d.get('usuario', ''),
            'status': d.get('status', ''),
            'atualizado_em': d.get('atualizado_em', ''),
        })
    return jsonify(result)


@app.route('/projeto/<doc_id>', methods=['GET'])
def carregar_projeto(doc_id):
    doc = db.collection('projetos').document(doc_id).get()
    if not doc.exists:
        return jsonify({'erro': 'Projeto não encontrado'}), 404
    d = doc.to_dict()
    return jsonify({'id': doc.id, 'nome': d.get('nome'), 'dados': d.get('dados', {}), 'status': d.get('status')})


@app.route('/projeto/<doc_id>', methods=['PUT'])
def atualizar_projeto(doc_id):
    data = request.json
    db.collection('projetos').document(doc_id).update({
        'nome': data.get('nome', ''),
        'dados': data.get('dados', {}),
        'status': data.get('status', 'rascunho'),
        'atualizado_em': datetime.now().isoformat(),
    })
    return jsonify({'ok': True})


@app.route('/projeto/<doc_id>', methods=['DELETE'])
def deletar_projeto(doc_id):
    db.collection('projetos').document(doc_id).delete()
    return jsonify({'ok': True})


# ── GERADOR DE SLIDES ──────────────────────────────────────────────────

_PPTX_DML = 'http://schemas.openxmlformats.org/drawingml/2006/main'


def _pptx_clone_shapes(target_slide, source_slide):
    tgt = target_slide.shapes._spTree
    src = source_slide.shapes._spTree
    for c in list(tgt)[2:]:
        tgt.remove(c)
    for c in list(src)[2:]:
        tgt.append(copy.deepcopy(c))


def _pptx_replace_text(shape, subs):
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        full = ''.join(r.text for r in para.runs)
        novo = full
        for old, new in subs.items():
            novo = _re.sub(_re.escape(old), new, novo, flags=_re.IGNORECASE)
        if novo != full and para.runs:
            para.runs[0].text = novo
            for r in para.runs[1:]:
                r.text = ''


def _pptx_fill_list(shape, fragment, lines):
    if not shape.has_text_frame or not lines:
        return False
    tf = shape.text_frame
    NS = _PPTX_DML
    for para in tf.paragraphs:
        if fragment.upper() in para.text.upper():
            p_elem = para._p
            parent = p_elem.getparent()
            idx = list(parent).index(p_elem)
            new_elems = []
            for line in lines:
                np = copy.deepcopy(p_elem)
                for r in np.findall(f'{{{NS}}}r'):
                    np.remove(r)
                orig_runs = p_elem.findall(f'{{{NS}}}r')
                if orig_runs:
                    new_r = copy.deepcopy(orig_runs[0])
                    t = new_r.find(f'{{{NS}}}t')
                    if t is not None:
                        t.text = line
                    np.append(new_r)
                new_elems.append(np)
            parent.remove(p_elem)
            for j, ne in enumerate(new_elems):
                parent.insert(idx + j, ne)
            return True
    return False


def _pptx_process_slide(slide, campos, imagens, modulo, modulo_descricao,
                         unidade, atividade, all_objs, all_enfases, all_etapas,
                         fase_info=None):
    sel = {}
    for campo in campos:
        tc = campo.get('tipo_campo', '')
        vals = campo.get('valores', [])
        if vals:
            sel[tc] = vals

    objs = sel.get('objetivos', all_objs)
    enfases = sel.get('enfases', all_enfases)
    etapas = sel.get('etapas', all_etapas)

    # Padrão genérico que captura (PUXAR...) com ou sem parêntese final
    _PAT = r'\(PUXAR[^)]*\)?'

    header_subs = {
        '(PUXAR MODULO)': modulo,
        '(PUXAR DESCRIÇÃO DA AULA)': atividade,
        '(PUXAR DESCRIÇÃO)': atividade,
    }

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        txt = '\n'.join(p.text for p in shape.text_frame.paragraphs).upper()

        if 'PUXAR MODULO' in txt or 'PUXAR DESCRIÇÃO' in txt:
            # Slides 3-6: cabeçalho "MÓDULO (PUXAR MODULO) – UNIDADE (PUXAR) – (PUXAR DESCRIÇÃO DA AULA)"
            _pptx_replace_text(shape, header_subs)
            for para in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in para.runs)
                novo = _re.sub(_PAT, unidade, full)
                if novo != full and para.runs:
                    para.runs[0].text = novo
                    for r in para.runs[1:]:
                        r.text = ''
        elif 'PUXAR OS OBJETIVOS' in txt:
            _pptx_fill_list(shape, 'PUXAR OS OBJETIVOS', objs)
        elif 'PUXAR NOME FASE' in txt:
            if fase_info:
                _pptx_replace_text(shape, {'PUXAR NOME FASE': fase_info.get('nome', '')})
        elif 'PUXAR ETAPAS' in txt:
            if fase_info:
                desc = fase_info.get('detalhe') or fase_info.get('nome', '')
                _pptx_fill_list(shape, 'PUXAR ETAPAS', [desc])
            else:
                _pptx_fill_list(shape, 'PUXAR ETAPAS', etapas)
        elif 'PUXAR A ENFASE' in txt:
            _pptx_fill_list(shape, 'PUXAR A ENFASE', enfases)
        elif 'PUXAR' in txt:
            # Slide 2 (VALORES): "Módulo (PUXAR): (PUXAR" e "Unidade (PUXAR) – (PUXAR)"
            for para in shape.text_frame.paragraphs:
                full = ''.join(r.text for r in para.runs)
                if 'PUXAR' not in full.upper():
                    continue
                novo = full
                tl = novo.lower()
                if 'módulo' in tl or 'modulo' in tl:
                    # 1º (PUXAR) → número do módulo, 2º → descrição do módulo
                    novo = _re.sub(_PAT, modulo, novo, count=1)
                    novo = _re.sub(_PAT, modulo_descricao, novo, count=1)
                elif 'unidade' in tl:
                    # 1º (PUXAR) → número da unidade, 2º → descrição da aula
                    novo = _re.sub(_PAT, unidade, novo, count=1)
                    novo = _re.sub(_PAT, atividade, novo, count=1)
                else:
                    novo = _re.sub(_PAT, unidade, novo)
                if novo != full and para.runs:
                    para.runs[0].text = novo
                    for r in para.runs[1:]:
                        r.text = ''

    for img_name in imagens:
        img_path = os.path.join(BASE_DIR, 'IMAGENS', img_name)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, PptxInches(1), PptxInches(3.5),
                                         width=PptxInches(9))
            except Exception:
                pass


@app.route('/gerar_slides', methods=['POST'])
def gerar_slides():
    form = request.form
    slide_config = json.loads(form.get('slide_config', '[]'))

    modulo = form.get('modulo', 'X')
    modulo_descricao = form.get('modulo_descricao', '')
    unidade = form.get('unidade', 'X')
    atividade = (form.get('atividade_descricao', '') or
                 form.get('descricao_aula', '') or
                 form.get('atividade_hidden', ''))

    all_objs = [v for v in form.getlist('obj_verbo') if v.strip()]
    all_enfases = [e for e in form.getlist('enfase_noun') if e.strip()]
    etapa_nomes = form.getlist('etapa_nome')
    etapa_outros = form.getlist('etapa_nome_outro')
    etapa_detalhes = form.getlist('etapa_detalhe')
    all_etapas = []
    all_etapas_info = []
    for i, nome in enumerate(etapa_nomes):
        n = nome.strip()
        if n == '__outro__' and i < len(etapa_outros):
            n = etapa_outros[i].strip()
        if n and n != '__outro__':
            detalhe = etapa_detalhes[i].strip() if i < len(etapa_detalhes) else ''
            all_etapas.append(n)
            all_etapas_info.append({'nome': n, 'detalhe': detalhe})

    template_path = os.path.join(BASE_DIR, 'MODELO SLIDE.pptx')
    prs_ref = PPTXPresentation(template_path)
    prs = PPTXPresentation(template_path)

    TIPO_REF = {
        'MISSÃO':     [prs_ref.slides[0], prs_ref.slides[6]],
        'VALORES':    [prs_ref.slides[1]],
        'OBJETIVOS':  [prs_ref.slides[2]],
        'FASES':      [prs_ref.slides[3]],
        'DEBRIEFING': [prs_ref.slides[4]],
        'ÊNFASE':     [prs_ref.slides[5]],
    }
    tipo_counts = {t: 0 for t in TIPO_REF}
    blank_layout = next(
        (l for l in prs.slide_layouts if l.name.lower() == 'blank'),
        prs.slide_layouts[-1]
    )

    expanded_config = []
    for cfg in slide_config:
        if cfg.get('tipo') == 'FASES' and all_etapas_info:
            for etapa in all_etapas_info:
                new_cfg = dict(cfg)
                new_cfg['_fase_info'] = etapa
                expanded_config.append(new_cfg)
        else:
            expanded_config.append(cfg)

    for i, cfg in enumerate(expanded_config):
        tipo = cfg.get('tipo', 'MISSÃO')
        campos = cfg.get('campos', [])
        imagens = cfg.get('imagens', [])
        fase_info = cfg.get('_fase_info')

        ref_list = TIPO_REF.get(tipo, TIPO_REF['MISSÃO'])
        cnt = tipo_counts.get(tipo, 0)
        ref_slide = ref_list[cnt % len(ref_list)]
        tipo_counts[tipo] = cnt + 1

        if i < 7:
            target = prs.slides[i]
        else:
            target = prs.slides.add_slide(blank_layout)

        _pptx_clone_shapes(target, ref_slide)
        _pptx_process_slide(target, campos, imagens, modulo, modulo_descricao,
                            unidade, atividade, all_objs, all_enfases, all_etapas,
                            fase_info=fase_info)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    fname = f"Slides_M{modulo}_U{unidade}.pptx"
    return send_file(buf, as_attachment=True, download_name=fname,
                     mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')


if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    app.run(debug=False, host='0.0.0.0', port=port)
